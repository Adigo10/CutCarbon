from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import copy

# Color palette
GREEN_DARK   = RGBColor(0x1A, 0x5C, 0x38)   # deep forest green
GREEN_MID    = RGBColor(0x2E, 0x7D, 0x4F)   # medium green
GREEN_LIGHT  = RGBColor(0xC8, 0xE6, 0xC9)   # pale green
TEAL         = RGBColor(0x00, 0x89, 0x7B)   # teal accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY         = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT    = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # completely blank


def add_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if transparency:
        shape.fill.fore_color.theme_color  # just access to confirm
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullet_box(slide, title, bullets, left, top, width, height,
                   title_size=16, bullet_size=13, bg_color=None,
                   title_color=GREEN_DARK, bullet_color=DARK_TEXT):
    if bg_color:
        add_rect(slide, left, top, width, height, bg_color)
    # title
    if title:
        add_text_box(slide, title, left + 0.15, top + 0.1, width - 0.3, 0.4,
                     font_size=title_size, bold=True, color=title_color)
    # bullets
    tb = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.55),
                                  Inches(width - 0.3), Inches(height - 0.65))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = bullet_color
        p.space_after = Pt(4)


# ─── SLIDE 1 – TITLE ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, GREEN_DARK)

# decorative accent bar
add_rect(s, 0, 5.8, 13.33, 1.7, GREEN_MID)
# teal strip
add_rect(s, 0, 3.0, 0.12, 2.8, TEAL)

add_text_box(s, "EventCarbon Co-Pilot", 0.5, 1.5, 12, 1.4,
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(s, "AI-Powered Carbon Footprint Management for Events",
             0.5, 3.1, 10, 0.8, font_size=22, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)
add_text_box(s, "Chat · Calculate · Compare · Reduce · Offset",
             0.5, 3.9, 10, 0.6, font_size=16, italic=True, color=TEAL, align=PP_ALIGN.LEFT)
add_text_box(s, "March 2026", 0.5, 6.2, 4, 0.5,
             font_size=14, color=GREEN_LIGHT, align=PP_ALIGN.LEFT)

# ─── SLIDE 2 – THE PROBLEM ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, LIGHT_GRAY)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_DARK)

add_text_box(s, "The Problem", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

stats = [
    ("🌍  Events are significant emitters", "Conferences, festivals & corporate events generate thousands of tCO₂e each — yet most organizers have no systematic way to measure or reduce them."),
    ("📋  Data is fragmented", "Emission data lives across caterers, venues, airlines, and logistics providers with no unified collection method."),
    ("🔢  Calculations are opaque", "Existing tools are spreadsheet-heavy, manual, and require sustainability expertise most organizers don't have."),
    ("💰  Carbon costs are rising", "Carbon taxes, SGX disclosure rules, EU CSRD & SBTi commitments mean accurate footprints are now a business necessity."),
]

for i, (heading, body) in enumerate(stats):
    col = i % 2
    row = i // 2
    lft = 0.3 + col * 6.5
    tp  = 1.3 + row * 2.7
    add_rect(s, lft, tp, 6.3, 2.5, WHITE)
    add_text_box(s, heading, lft + 0.2, tp + 0.15, 5.9, 0.5,
                 font_size=14, bold=True, color=GREEN_DARK)
    add_text_box(s, body, lft + 0.2, tp + 0.7, 5.9, 1.6,
                 font_size=12, color=GRAY, wrap=True)

# ─── SLIDE 3 – SOLUTION OVERVIEW ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_MID)

add_text_box(s, "Our Solution", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

add_text_box(s,
    "EventCarbon Co-Pilot is a chat-first, AI-powered platform that guides event organizers "
    "from raw natural-language descriptions to verified GHG reports — in minutes, not weeks.",
    0.5, 1.2, 12.3, 1.0, font_size=16, color=DARK_TEXT, wrap=True)

pillars = [
    ("Chat-First\nIntake", "Describe your event in plain English. The AI extracts structured activity data automatically."),
    ("Deterministic\nCalculations", "GHG Protocol-compliant engine. No LLM guesswork — every tCO₂e is traceable to a factor."),
    ("Multi-Scenario\nComparison", "Clone & tweak scenarios side-by-side: travel mode shifts, menu changes, venue swaps."),
    ("Live Factor\nUpdates", "Web agents refresh emission factors from DEFRA, ICAO, IEA, and NEA automatically."),
    ("Financial\n& Compliance", "Carbon tax savings, green incentives, and compliance scoring for SBTi, SGX, EU CSRD."),
]

for i, (title, desc) in enumerate(pillars):
    lft = 0.3 + i * 2.55
    add_rect(s, lft, 2.4, 2.4, 4.4, GREEN_DARK)
    add_rect(s, lft, 2.4, 2.4, 0.08, TEAL)  # top accent
    add_text_box(s, title, lft + 0.1, 2.5, 2.2, 0.9,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, desc, lft + 0.1, 3.5, 2.2, 3.1,
                 font_size=11, color=GREEN_LIGHT, wrap=True, align=PP_ALIGN.CENTER)

# ─── SLIDE 4 – USER ROLES ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, LIGHT_GRAY)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_DARK)

add_text_box(s, "Who Uses It", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

roles = [
    ("🎪  Event Organizer", "Primary user", [
        "Defines event parameters via chat",
        "Collects supplier data",
        "Runs scenario comparisons",
        "Generates GHG reports for sponsors",
    ]),
    ("🌿  Sustainability Lead", "Expert user", [
        "Fine-tunes emission factor assumptions",
        "Uploads precise supplier data",
        "Validates methodology boundaries",
        "Benchmarks against SBTi / ISO 20121",
    ]),
    ("📊  Sponsor / Client", "Read-only stakeholder", [
        "Views footprint + reduction + offset report",
        "Reviews compliance scoring",
        "Compares year-over-year progress",
        "Exports PDF / CSV for ESG disclosure",
    ]),
]

for i, (role, subtitle, bullets) in enumerate(roles):
    lft = 0.4 + i * 4.3
    add_rect(s, lft, 1.3, 4.1, 5.8, WHITE)
    add_rect(s, lft, 1.3, 4.1, 0.08, TEAL)
    add_text_box(s, role, lft + 0.2, 1.4, 3.7, 0.5,
                 font_size=16, bold=True, color=GREEN_DARK)
    add_text_box(s, subtitle, lft + 0.2, 1.95, 3.7, 0.35,
                 font_size=12, italic=True, color=GRAY)
    tb = s.shapes.add_textbox(Inches(lft + 0.2), Inches(2.4), Inches(3.7), Inches(4.4))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for j, b in enumerate(bullets):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

# ─── SLIDE 5 – KEY USER FLOWS ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_MID)

add_text_box(s, "Key User Flows", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

flows = [
    ("1  Chat-First Intake",
     '"We\'re planning a 3-day tech conference in Singapore, 800 attendees, hybrid, mostly APAC audience."',
     "AI asks clarifying questions → generates initial scenario using proxy factors → all assumptions clearly flagged."),
    ("2  Multi-Scenario Comparison",
     '"What if we move it to Bangkok and push for more rail travel?"',
     "Agent clones baseline, tweaks parameters → shows Scenario A vs B vs C: total tCO₂e, per-attendee, category breakdown."),
    ("3  Supplier Data Collection",
     '"Ask our caterer and AV vendor for precise data."',
     "Platform generates shareable mini-forms. Suppliers submit kWh, fuel, kg, menu data → chatbot ingests & recalculates."),
    ("4  Reduction & Offset Planning",
     '"How can we cut 30% of emissions?"',
     "Agent suggests measure sets (travel mode shift, menu change, logistics consolidation) and estimates remaining offset need."),
]

for i, (title, quote, desc) in enumerate(flows):
    col = i % 2
    row = i // 2
    lft = 0.3 + col * 6.5
    tp  = 1.25 + row * 2.9
    add_rect(s, lft, tp, 6.3, 2.7, LIGHT_GRAY)
    add_rect(s, lft, tp, 0.08, 2.7, TEAL)
    add_text_box(s, title, lft + 0.25, tp + 0.1, 5.9, 0.45,
                 font_size=14, bold=True, color=GREEN_DARK)
    add_text_box(s, quote, lft + 0.25, tp + 0.6, 5.9, 0.7,
                 font_size=11, italic=True, color=TEAL, wrap=True)
    add_text_box(s, desc, lft + 0.25, tp + 1.35, 5.9, 1.2,
                 font_size=11, color=GRAY, wrap=True)

# ─── SLIDE 6 – ARCHITECTURE ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, LIGHT_GRAY)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_DARK)

add_text_box(s, "Architecture Overview", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

layers = [
    ("Frontend  (React 19 + TypeScript + Vite)",
     ["Chat pane — natural-language event intake & Q&A",
      "Structured side-panel — auto-filled forms (attendees, travel, venue, catering, waste)",
      "Scenario dashboard — cards, stacked bar charts, comparison table",
      "Financial panel — carbon tax savings, incentives, compliance scores"]),
    ("Orchestrator API  (FastAPI async)",
     ["POST /api/chat — OpenAI function calling → structured EventScenarioInput",
      "CRUD + clone + compare — /api/scenarios",
      "Financial analysis — /api/financial  (tax, incentives, compliance by region)",
      "TinyFish agent triggers — /api/agents"]),
    ("Services Layer",
     ["emissions_engine.py — deterministic GHG Protocol calculations",
      "financial_engine.py — carbon tax, incentives, compliance scoring",
      "openai_service.py — NL→structured extraction + recommendations",
      "tinyfish_agent.py — 6 headless browser agents for live factor refresh"]),
    ("Data / Storage",
     ["SQLite (dev) / PostgreSQL (prod) — async SQLAlchemy ORM",
      "emission_factors.json — factor catalog with provenance & methodology tags",
      "tax_incentives.json — rates, exchange rates, incentives by region (SG/EU/UK/AU/USA)"]),
]

colors = [GREEN_DARK, GREEN_MID, TEAL, RGBColor(0x00, 0x6D, 0x77)]
for i, (title, pts) in enumerate(layers):
    lft = 0.25
    tp  = 1.2 + i * 1.5
    add_rect(s, lft, tp, 12.8, 1.4, colors[i])
    add_text_box(s, title, lft + 0.15, tp + 0.05, 4.0, 0.45,
                 font_size=13, bold=True, color=WHITE)
    text = "   ·  " + "   ·  ".join(pts)
    add_text_box(s, text, lft + 4.2, tp + 0.05, 8.4, 1.25,
                 font_size=11, color=WHITE, wrap=True)

# ─── SLIDE 7 – EMISSIONS ENGINE ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_MID)

add_text_box(s, "Emissions Engine", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

add_text_box(s,
    "100% deterministic — the LLM never touches calculations. "
    "Every tCO₂e is traceable to a published emission factor.",
    0.4, 1.2, 12.5, 0.6, font_size=14, italic=True, color=DARK_TEXT, wrap=True)

modules = [
    ("✈  Travel",       "participants × distance × factor\n(per mode + class: economy/business/first, car, rail, ferry)"),
    ("⚡  Venue Energy",  "kWh × grid factor  (or proxy per m²/day in basic mode)"),
    ("🏨  Accommodation", "nights × room type × factor"),
    ("🍽  Catering",      "meals × menu type (meat / vegetarian / vegan) × factor"),
    ("♻  Materials & Waste", "mass or units × waste / emission factor"),
]

for i, (mod, formula) in enumerate(modules):
    col = i % 3
    row = i // 3
    lft = 0.3 + col * 4.3
    tp  = 2.0 + row * 2.4
    add_rect(s, lft, tp, 4.1, 2.2, GREEN_DARK)
    add_rect(s, lft, tp, 4.1, 0.07, TEAL)
    add_text_box(s, mod, lft + 0.15, tp + 0.1, 3.8, 0.5,
                 font_size=14, bold=True, color=WHITE)
    add_text_box(s, formula, lft + 0.15, tp + 0.7, 3.8, 1.3,
                 font_size=11, color=GREEN_LIGHT, wrap=True)

# modes label
add_rect(s, 9.0, 2.0, 4.1, 2.2, LIGHT_GRAY)
add_text_box(s, "Two Calculation Modes", 9.15, 2.1, 3.8, 0.5,
             font_size=13, bold=True, color=GREEN_DARK)
add_text_box(s,
    "Basic — proxy-heavy, great for early planning. Estimates clearly flagged.\n\n"
    "Advanced — detailed supplier data, scope 1/2/3 fully resolved, audit-ready.",
    9.15, 2.65, 3.8, 1.4, font_size=11, color=DARK_TEXT, wrap=True)

# ─── SLIDE 8 – AI & WEB AGENTS ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, LIGHT_GRAY)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_DARK)

add_text_box(s, "AI Layer & Web Agents", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

# Left column – OpenAI
add_rect(s, 0.3, 1.25, 6.1, 5.9, WHITE)
add_rect(s, 0.3, 1.25, 6.1, 0.08, TEAL)
add_text_box(s, "🤖  OpenAI Chat Service", 0.5, 1.35, 5.7, 0.5,
             font_size=16, bold=True, color=GREEN_DARK)
chat_pts = [
    "NL → structured activity data\n  \"70% fly EU, 30% local by MRT\" → typed breakdown",
    "Sensible defaults when data is missing\n  (clearly flagged as estimates)",
    "Reduction & offset recommendations\n  with methodology explanations",
    "Three chatbot modes:\n  • Data collection — fills activity items progressively\n  • Explanation — answers 'why is this number high?'\n  • What-if — 'clone baseline, switch meals to vegan'",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.7), Inches(5.0))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
for j, b in enumerate(chat_pts):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(10)

# Right column – TinyFish
add_rect(s, 6.9, 1.25, 6.1, 5.9, WHITE)
add_rect(s, 6.9, 1.25, 6.1, 0.08, TEAL)
add_text_box(s, "🌐  TinyFish Web Agents (×6)", 7.1, 1.35, 5.7, 0.5,
             font_size=16, bold=True, color=GREEN_DARK)
agent_pts = [
    "EMA — Singapore grid emission factors",
    "DEFRA — UK/EU emission factors",
    "NEA — National Environment Agency factors",
    "Ember Climate — live grid carbon intensity",
    "ICAO — aviation emission factors",
    "Our World in Data — global energy mix",
]
add_text_box(s, "Agents run on schedule and merge results into emission_factors.json with full provenance (source URL, fetch time, run_id).",
             7.1, 2.0, 5.7, 0.8, font_size=12, color=GRAY, wrap=True)
tb2 = s.shapes.add_textbox(Inches(7.1), Inches(2.9), Inches(5.7), Inches(3.8))
tb2.word_wrap = True
tf2 = tb2.text_frame
tf2.word_wrap = True
for j, b in enumerate(agent_pts):
    p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(8)

# ─── SLIDE 9 – DATA MODEL ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_MID)

add_text_box(s, "Core Data Model", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

entities = [
    ("Event",          "id · name · location · dates · type · organizer_id"),
    ("Scenario",       "id · event_id · name · description · assumptions (JSON) · mode (basic/advanced)"),
    ("ActivityGroup",  "type (travel/venue_energy/accommodation/catering/materials/waste) · scope (1/2/3) · scenario_id"),
    ("ActivityItem",   "group_id · quantity · unit · emission_factor_id · source (user_input / chat_inferred / supplier_form / default_proxy)"),
    ("EmissionFactor", "id · category · region · value · unit · source_url · last_updated · methodology_tag (IPCC2013 / GHGProtocol)"),
    ("FinancialReport","id · scenario_id · region · carbon_tax_savings · incentives (JSON) · compliance_score"),
]

for i, (entity, fields) in enumerate(entities):
    col = i % 2
    row = i // 2
    lft = 0.3 + col * 6.5
    tp  = 1.25 + row * 1.9
    add_rect(s, lft, tp, 6.3, 1.75, LIGHT_GRAY)
    add_rect(s, lft, tp, 0.08, 1.75, GREEN_MID)
    add_text_box(s, entity, lft + 0.25, tp + 0.1, 5.8, 0.4,
                 font_size=14, bold=True, color=GREEN_DARK)
    add_text_box(s, fields, lft + 0.25, tp + 0.55, 5.8, 1.0,
                 font_size=11, color=GRAY, wrap=True)

# ─── SLIDE 10 – FINANCIAL & COMPLIANCE ───────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, LIGHT_GRAY)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_DARK)

add_text_box(s, "Financial & Compliance Engine", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

regions = [
    ("🇸🇬 Singapore", "Carbon tax S$25/tCO₂e\nNEA Green Mark incentives\nSGX sustainability disclosure"),
    ("🇪🇺 European Union", "ETS price ~€65/tCO₂e\nEU Green Deal incentives\nEU CSRD mandatory reporting"),
    ("🇬🇧 United Kingdom", "UK ETS + carbon price support\nBBSR Green Building Fund\nISO 20121 alignment"),
    ("🇦🇺 Australia", "Safeguard Mechanism\nERF incentive programs\nASRS climate disclosures"),
    ("🇺🇸 United States", "State-level carbon markets\nIRA clean energy credits\nSEC climate rules"),
]

for i, (region, details) in enumerate(regions):
    lft = 0.25 + i * 2.55
    add_rect(s, lft, 1.3, 2.4, 3.2, GREEN_DARK)
    add_rect(s, lft, 1.3, 2.4, 0.07, TEAL)
    add_text_box(s, region, lft + 0.1, 1.4, 2.2, 0.5,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, details, lft + 0.1, 2.0, 2.2, 2.3,
                 font_size=10, color=GREEN_LIGHT, wrap=True, align=PP_ALIGN.CENTER)

add_rect(s, 0.25, 4.8, 12.8, 2.4, WHITE)
add_text_box(s, "Compliance Frameworks Scored Automatically",
             0.5, 4.9, 12.2, 0.5, font_size=15, bold=True, color=GREEN_DARK)
frameworks = "GHG Protocol (Scope 1/2/3)   ·   ISO 20121 (Sustainable Events)   ·   SBTi (Science-Based Targets)   ·   SGX Sustainability Reporting   ·   EU CSRD"
add_text_box(s, frameworks, 0.5, 5.5, 12.2, 0.8,
             font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ─── SLIDE 11 – MVP vs STRETCH ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.33, 1.1, GREEN_MID)

add_text_box(s, "Roadmap: MVP → Full Platform", 0.4, 0.15, 12, 0.75,
             font_size=32, bold=True, color=WHITE)

# MVP column
add_rect(s, 0.3, 1.25, 5.9, 5.9, GREEN_DARK)
add_rect(s, 0.3, 1.25, 5.9, 0.07, TEAL)
add_text_box(s, "MVP  (Hackathon — 24–36h)", 0.5, 1.35, 5.5, 0.5,
             font_size=16, bold=True, color=WHITE)
mvp_pts = [
    "Singapore / EU generic emission factors",
    "Categories: Travel + Venue Energy + Catering + simple Waste proxy",
    "Chat intake → structured scenario → deterministic calculation",
    "Scenario comparison chart (stacked bar)",
    "Carbon tax savings estimate (SG & EU)",
    "Basic compliance score",
    "React dashboard + FastAPI backend",
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(4.9))
tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
for j, b in enumerate(mvp_pts):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.text = f"✓  {b}"; p.font.size = Pt(12); p.font.color.rgb = GREEN_LIGHT; p.space_after = Pt(8)

# Stretch column
add_rect(s, 6.9, 1.25, 6.1, 5.9, LIGHT_GRAY)
add_rect(s, 6.9, 1.25, 6.1, 0.07, TEAL)
add_text_box(s, "Stretch Goals", 7.1, 1.35, 5.7, 0.5,
             font_size=16, bold=True, color=GREEN_DARK)
stretch_pts = [
    "Supplier mini-forms with ingestion pipeline",
    "Multi-event history & year-over-year progress view",
    "Automatic offset suggestion (links to offset providers via web agents)",
    "Advanced mode with full scope 1/2/3 breakdown",
    "Scenario versioning & audit trail",
    "PDF/CSV export for ESG disclosure",
    "Real-time TinyFish factor refresh on demand",
    "Multi-region carbon tax & incentive matching",
]
tb2 = s.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.7), Inches(4.9))
tb2.word_wrap = True; tf2 = tb2.text_frame; tf2.word_wrap = True
for j, b in enumerate(stretch_pts):
    p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
    p.text = f"→  {b}"; p.font.size = Pt(12); p.font.color.rgb = DARK_TEXT; p.space_after = Pt(7)

# ─── SLIDE 12 – THANK YOU ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(s, GREEN_DARK)
add_rect(s, 0, 3.5, 13.33, 0.07, TEAL)
add_rect(s, 0, 0, 0.12, 7.5, TEAL)

add_text_box(s, "Thank You", 0.5, 1.2, 12, 1.2,
             font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s,
    "EventCarbon Co-Pilot  ·  Chat · Calculate · Compare · Reduce · Offset",
    0.5, 2.7, 12, 0.6, font_size=18, italic=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text_box(s,
    "Questions?",
    0.5, 4.2, 12, 0.7, font_size=24, bold=True, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(s,
    "Built with FastAPI · React · OpenAI · TinyFish Agents · GHG Protocol",
    0.5, 5.5, 12, 0.5, font_size=13, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

# ─── SAVE ─────────────────────────────────────────────────────────────────────
out = "EventCarbon_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
